# -*- coding: utf-8 -*-
"""One-off builder: cleaning_operations_demo_presentation.pptx from scshots/*.png"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "scshots"
OUT_PPTX = ROOT / "cleaning_operations_demo_presentation.pptx"
OUT_PDF = ROOT / "cleaning_operations_demo_presentation.pdf"

SCENES = [
    (
        "Manager-side operational evidence overview",
        "Single-task visibility of QR entry, late check-in, timestamps, GPS, duration, and photo evidence.",
        1,
    ),
    (
        "Assignment, site, and planned visit overview",
        "Manager can see the visit title, customer site, planned window, and assigned cleaner.",
        2,
    ),
    (
        "Cleaner portal homepage and entry point",
        "Cleaner accesses the workflow from a simple portal homepage through My Visits.",
        3,
    ),
    (
        "Cleaner work queue / My Visits list",
        "Cleaner sees only assigned visits with clear operational statuses.",
        4,
    ),
    (
        "Visit detail before starting execution",
        "Visit starts from a simple portal page with Start Visit and photo evidence sections.",
        5,
    ),
    (
        "Visit started with graceful GPS fallback",
        "Check-in is recorded even when GPS cannot be captured.",
        6,
    ),
    (
        "Before and after photo evidence captured",
        "Cleaner uploads visual evidence directly from the portal.",
        7,
    ),
    (
        "Visit completed with check-out and duration",
        "System records check-out and computes total visit duration.",
        8,
    ),
    (
        "Late check-in visibility for the manager",
        "Manager can clearly identify delayed visit starts versus the planned schedule.",
        9,
    ),
    (
        "QR-based entry URL for the visit",
        "Each visit can be opened through a lightweight URL-based QR entry flow.",
        10,
    ),
    (
        "Printable portal visit summary report",
        "Manager can export a clean single-visit summary for review and presentation.",
        11,
    ),
]

CLOSING_BULLETS = [
    "Manager-side visit assignment and oversight",
    "Cleaner portal execution flow",
    "Start and end timestamps",
    "Optional GPS at check-in with graceful fallback",
    "Before and after photo evidence",
    "Late check-in visibility",
    "QR-based visit entry",
    "Printable visit summary report",
]


def _set_textbox(tb, text: str, size_pt: float, bold: bool = False, color: RGBColor | None = None):
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = "Calibri"
    if color:
        run.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def _add_scene_slide(prs, title: str, subtitle: str, image_path: Path):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    box_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.65))
    _set_textbox(box_t, title, 22, bold=True, color=RGBColor(0x20, 0x20, 0x20))

    # Subtitle
    box_s = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.75))
    _set_textbox(box_s, subtitle, 14, bold=False, color=RGBColor(0x44, 0x44, 0x44))

    # Image scaled to fit content area (16:9 slide)
    img = Image.open(image_path)
    w_px, h_px = img.size
    aspect = h_px / w_px

    max_w_in = 12.3
    max_h_in = 5.85
    w_in = max_w_in
    h_in = w_in * aspect
    if h_in > max_h_in:
        h_in = max_h_in
        w_in = h_in / aspect

    left = (float(prs.slide_width) / 914400 - w_in) / 2
    top = 1.75
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(w_in), height=Inches(h_in))


def main() -> None:
    missing = []
    for i in range(1, 12):
        p = SHOTS / f"{i}.png"
        if not p.is_file():
            missing.append(str(p))
    if missing:
        raise FileNotFoundError("Missing screenshots:\n" + "\n".join(missing))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1) Title slide
    tslide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = tslide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.8), Inches(1.2))
    _set_textbox(tbox, "Cleaning Operations Workflow Demo", 36, bold=True, color=RGBColor(0x1A, 0x37, 0x56))
    tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    sbox = tslide.shapes.add_textbox(Inches(0.8), Inches(3.55), Inches(11.8), Inches(1.0))
    _set_textbox(
        sbox,
        "Odoo 19 Field Service + Portal Cleaner Execution Flow",
        18,
        bold=False,
        color=RGBColor(0x44, 0x44, 0x44),
    )
    sbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 2–12) Scene slides
    for title, subtitle, num in SCENES:
        _add_scene_slide(prs, f"Scene {num} — {title}", subtitle, SHOTS / f"{num}.png")

    # 13) Closing
    cslide = prs.slides.add_slide(prs.slide_layouts[6])
    ctitle = cslide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.7))
    _set_textbox(ctitle, "Summary", 28, bold=True, color=RGBColor(0x20, 0x20, 0x20))

    bullets = cslide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(11.9), Inches(5.8))
    tf = bullets.text_frame
    tf.clear()
    tf.paragraphs[0].text = CLOSING_BULLETS[0]
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].space_after = Pt(10)
    for line in CLOSING_BULLETS[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(10)

    prs.save(OUT_PPTX)
    print(f"Saved: {OUT_PPTX}")


if __name__ == "__main__":
    main()
